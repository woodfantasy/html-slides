#!/usr/bin/env python3
"""
html-slides PPT Extractor
Extracts text, images, and speaker notes from .pptx files
to generate an outline for html-slides content filling.

Usage:
    python extract-pptx.py input.pptx output_dir/

Output:
    output_dir/
    ├── outline.md          # Markdown outline with all slide text
    ├── images/             # Extracted images (renamed)
    │   ├── 01-image1.png
    │   ├── 02-image1.jpg
    │   └── ...
    └── notes.md            # Speaker notes per slide

Requirements: python-pptx (pip install python-pptx)
"""

import sys
import os
import re
from pathlib import Path

def check_dependencies():
    try:
        from pptx import Presentation
        return True
    except ImportError:
        print("❌ Missing dependency: python-pptx")
        print("   Install: pip install python-pptx")
        return False


def sanitize_filename(name):
    """Remove invalid characters from filename."""
    return re.sub(r'[^\w\-.]', '_', name)


def extract_pptx(input_path, output_dir):
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.enum.text import PP_ALIGN

    if not os.path.exists(input_path):
        print(f"❌ File not found: {input_path}")
        sys.exit(1)

    # Create output directories
    output_path = Path(output_dir)
    images_path = output_path / "images"
    output_path.mkdir(parents=True, exist_ok=True)
    images_path.mkdir(exist_ok=True)

    prs = Presentation(input_path)
    total_slides = len(prs.slides)

    print(f"📊 Extracting: {input_path}")
    print(f"   Slides: {total_slides}")
    print(f"   Size: {prs.slide_width.emu}x{prs.slide_height.emu} EMU")
    print()

    outline_lines = ["# PPT Outline\n"]
    outline_lines.append(f"> Source: `{os.path.basename(input_path)}`\n")
    outline_lines.append(f"> Slides: {total_slides}\n\n---\n")

    notes_lines = ["# Speaker Notes\n"]
    notes_lines.append(f"> Source: `{os.path.basename(input_path)}`\n\n---\n")

    image_count = 0

    for slide_idx, slide in enumerate(prs.slides, 1):
        slide_num = str(slide_idx).zfill(2)

        # Extract layout name
        layout_name = slide.slide_layout.name if slide.slide_layout else "Unknown"

        outline_lines.append(f"\n## Slide {slide_num} [{layout_name}]\n")

        # Extract text from all shapes
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        # Detect title vs body by font size
                        is_title = False
                        for run in para.runs:
                            if run.font.size and run.font.size > Pt(24):
                                is_title = True
                                break
                            if run.font.bold:
                                is_title = True
                                break

                        if is_title:
                            texts.append(f"**{text}**")
                        else:
                            texts.append(text)

            # Extract images
            if shape.shape_type == 13:  # Picture
                image_count += 1
                image = shape.image
                ext = image.content_type.split('/')[-1]
                if ext == 'jpeg': ext = 'jpg'
                img_name = f"{slide_num}-image{image_count}.{ext}"
                img_path = images_path / img_name
                with open(img_path, 'wb') as f:
                    f.write(image.blob)
                texts.append(f"![Image]({img_name})")
                print(f"  📷 Slide {slide_num}: Extracted {img_name}")

            # Extract tables
            if shape.has_table:
                table = shape.table
                texts.append("\n| " + " | ".join(cell.text for cell in table.rows[0].cells) + " |")
                texts.append("| " + " | ".join("---" for _ in table.rows[0].cells) + " |")
                for row in table.rows[1:]:
                    texts.append("| " + " | ".join(cell.text for cell in row.cells) + " |")
                texts.append("")

        for text in texts:
            outline_lines.append(f"- {text}\n")

        # Extract speaker notes
        if slide.has_notes_slide:
            notes_text = slide.notes_slide.notes_text_frame.text.strip()
            if notes_text:
                notes_lines.append(f"\n### Slide {slide_num}\n")
                notes_lines.append(f"{notes_text}\n")

    # Write outline
    outline_path = output_path / "outline.md"
    with open(outline_path, 'w', encoding='utf-8') as f:
        f.write('\n'.join(outline_lines))

    # Write notes
    notes_path = output_path / "notes.md"
    with open(notes_path, 'w', encoding='utf-8') as f:
        f.write('\n'.join(notes_lines))

    print(f"\n✅ Extraction complete!")
    print(f"   📝 Outline: {outline_path}")
    print(f"   🗒️  Notes:   {notes_path}")
    print(f"   🖼️  Images:  {image_count} extracted to {images_path}/")
    print()
    print("Next steps:")
    print("  1. Review outline.md for content structure")
    print("  2. Choose a style and copy the template")
    print("  3. Fill in content using layouts from references/layouts.md")
    print("  4. Move images/ next to your index.html")


if __name__ == "__main__":
    if len(sys.argv) < 3:
        print("Usage: python extract-pptx.py <input.pptx> <output_dir/>")
        sys.exit(1)

    if not check_dependencies():
        sys.exit(1)

    extract_pptx(sys.argv[1], sys.argv[2])
